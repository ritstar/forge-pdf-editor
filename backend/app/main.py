from __future__ import annotations

import logging
import os
import re
import secrets
import shutil
import subprocess
import tempfile
import time
from pathlib import Path
from threading import Lock
from typing import Optional
from urllib.parse import urlparse

import fitz
import pikepdf
from fastapi import FastAPI, File, Form, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from openpyxl import Workbook
from pydantic import BaseModel
from pypdf import PdfReader, PdfWriter
from starlette.background import BackgroundTask
from weasyprint import HTML, default_url_fetcher
from docx import Document
from pptx import Presentation


APP_TITLE = "Forge PDF Tools API"
APP_VERSION = "0.1.0"
API_KEY_HEADER = "x-api-key"
API_KEY_ENV = "PDF_TOOLS_API_SECRET"
MAX_HTML_INPUT_BYTES = 200_000
RATE_LIMIT_MAX_REQUESTS = 20
RATE_LIMIT_WINDOW_SECONDS = 60

logger = logging.getLogger(__name__)
_rate_limit_lock = Lock()
_rate_limit_state: dict[str, list[float]] = {}


def _allowed_origins() -> list[str]:
    raw = os.getenv("CORS_ORIGINS", "http://localhost:3000")
    return [x.strip() for x in raw.split(",") if x.strip()]


def _expected_api_key() -> str:
    value = os.getenv(API_KEY_ENV, "").strip()
    if not value:
        logger.error("Missing required server env for backend API authentication.")
        raise HTTPException(status_code=500, detail="Server misconfiguration.")
    return value


def _require_api_key(request: Request) -> None:
    provided = request.headers.get(API_KEY_HEADER, "")
    expected = _expected_api_key()
    if not secrets.compare_digest(provided, expected):
        raise HTTPException(status_code=401, detail="Unauthorized")


def _client_ip(request: Request) -> str:
    forwarded = request.headers.get("x-forwarded-for")
    if forwarded:
        return forwarded.split(",")[0].strip()
    return request.headers.get("x-real-ip") or "unknown"


def _enforce_rate_limit(request: Request) -> None:
    current = time.monotonic()
    key = _client_ip(request)

    with _rate_limit_lock:
        history = _rate_limit_state.get(key, [])
        history = [item for item in history if current - item < RATE_LIMIT_WINDOW_SECONDS]

        if len(history) >= RATE_LIMIT_MAX_REQUESTS:
            _rate_limit_state[key] = history
            raise HTTPException(status_code=429, detail="Too many requests.")

        history.append(current)
        _rate_limit_state[key] = history


app = FastAPI(title=APP_TITLE, version=APP_VERSION)

app.add_middleware(
    CORSMiddleware,
    allow_origins=_allowed_origins(),
    allow_methods=["*"],
    allow_headers=["*"],
)


class HealthResponse(BaseModel):
    status: str
    app: str
    version: str


@app.get("/health", response_model=HealthResponse)
def health() -> HealthResponse:
    return HealthResponse(status="ok", app=APP_TITLE, version=APP_VERSION)


def _safe_stem(name: str, default: str = "output") -> str:
    stem = Path(name).stem or default
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "-", stem).strip("-")
    return cleaned or default


def _temp_dir() -> Path:
    return Path(tempfile.mkdtemp(prefix="forge-tools-"))


def _save_upload(upload: UploadFile, target: Path) -> Path:
    with target.open("wb") as f:
        shutil.copyfileobj(upload.file, f)
    return target


def _as_download(path: Path, filename: str, tmpdir: Path) -> FileResponse:
    return FileResponse(
        path,
        media_type="application/octet-stream",
        filename=filename,
        background=BackgroundTask(lambda: shutil.rmtree(tmpdir, ignore_errors=True)),
    )


def _safe_html_fetcher(url: str, *args, **kwargs):
    parsed = urlparse(url)
    if parsed.scheme != "data":
        raise ValueError("External resources are not allowed for HTML to PDF conversion.")
    return default_url_fetcher(url, *args, **kwargs)


def _pdf_to_word(input_path: Path, output_path: Path) -> None:
    reader = PdfReader(str(input_path))
    doc = Document()

    for i, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        doc.add_heading(f"Page {i + 1}", level=2)
        doc.add_paragraph(text if text else "")
        if i != len(reader.pages) - 1:
            doc.add_page_break()

    doc.save(str(output_path))


def _pdf_to_excel(input_path: Path, output_path: Path) -> None:
    reader = PdfReader(str(input_path))
    wb = Workbook()
    ws = wb.active
    ws.title = "PDF Text"
    ws.append(["page", "line", "text"])

    for page_idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        lines = [ln for ln in text.splitlines() if ln.strip()]
        if not lines:
            ws.append([page_idx, 1, ""])
            continue

        for line_idx, line in enumerate(lines, start=1):
            ws.append([page_idx, line_idx, line])

    wb.save(str(output_path))


def _pdf_to_powerpoint(input_path: Path, output_path: Path, tmpdir: Path) -> None:
    pdf = fitz.open(str(input_path))
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for i, page in enumerate(pdf):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        img_path = tmpdir / f"page-{i + 1}.png"
        pix.save(str(img_path))

        slide = prs.slides.add_slide(blank_layout)
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Fill entire slide with rendered page.
        slide.shapes.add_picture(str(img_path), 0, 0, width=slide_w, height=slide_h)

    prs.save(str(output_path))


def _libreoffice_convert(input_path: Path, output_ext: str, out_dir: Path) -> Path:
    cmd = [
        "soffice",
        "--headless",
        "--convert-to",
        output_ext,
        "--outdir",
        str(out_dir),
        str(input_path),
    ]

    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if proc.returncode != 0:
        logger.error("LibreOffice conversion failed: %s", (proc.stderr or proc.stdout).strip())
        raise HTTPException(status_code=500, detail="Document conversion failed.")

    out_path = out_dir / f"{input_path.stem}.{output_ext}"
    if not out_path.exists():
        raise HTTPException(status_code=500, detail="Converted output file not found.")

    return out_path


def _convert_to_pdfa(input_path: Path, output_path: Path) -> None:
    cmd = [
        "gs",
        "-dPDFA=2",
        "-dBATCH",
        "-dNOPAUSE",
        "-dNOOUTERSAVE",
        "-sProcessColorModel=DeviceRGB",
        "-sDEVICE=pdfwrite",
        "-sPDFACompatibilityPolicy=1",
        f"-sOutputFile={output_path}",
        str(input_path),
    ]

    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if proc.returncode != 0:
        logger.error("PDF/A conversion failed: %s", (proc.stderr or proc.stdout).strip())
        raise HTTPException(status_code=500, detail="PDF/A conversion failed.")


@app.post("/tools/{tool_id}")
async def run_tool(
    request: Request,
    tool_id: str,
    file: Optional[UploadFile] = File(default=None),
    password: str = Form(default=""),
    owner_password: str = Form(default=""),
    compression_level: str = Form(default="medium"),
    html_content: str = Form(default=""),
):
    _require_api_key(request)
    _enforce_rate_limit(request)
    tmpdir = _temp_dir()

    try:
        if tool_id == "html-to-pdf":
            if file is None and not html_content.strip():
                raise HTTPException(status_code=400, detail="Provide either an HTML file or html_content.")

            if file is not None:
                html_text = (await file.read()).decode("utf-8", errors="replace")
                base_name = _safe_stem(file.filename or "html")
            else:
                html_text = html_content
                base_name = "html"

            if len(html_text.encode("utf-8")) > MAX_HTML_INPUT_BYTES:
                raise HTTPException(status_code=400, detail="HTML input is too large.")

            output_path = tmpdir / f"{base_name}.pdf"
            HTML(string=html_text, url_fetcher=_safe_html_fetcher).write_pdf(str(output_path))
            return _as_download(output_path, output_path.name, tmpdir)

        if file is None:
            raise HTTPException(status_code=400, detail="A file is required for this tool.")

        filename = file.filename or "input"
        input_path = tmpdir / filename
        _save_upload(file, input_path)

        stem = _safe_stem(filename)

        if tool_id == "compress-pdf":
            output_path = tmpdir / f"{stem}_compressed.pdf"
            with pikepdf.open(str(input_path), allow_overwriting_input=True) as pdf:
                decode_level = {
                    "low": pikepdf.StreamDecodeLevel.all,
                    "medium": pikepdf.StreamDecodeLevel.generalized,
                    "high": pikepdf.StreamDecodeLevel.specialized,
                }.get(compression_level, pikepdf.StreamDecodeLevel.generalized)
                pdf.save(
                    str(output_path),
                    compress_streams=True,
                    stream_decode_level=decode_level,
                    object_stream_mode=pikepdf.ObjectStreamMode.generate,
                    recompress_flate=True,
                )
            return _as_download(output_path, output_path.name, tmpdir)

        if tool_id == "pdf-to-word":
            output_path = tmpdir / f"{stem}.docx"
            _pdf_to_word(input_path, output_path)
            return _as_download(output_path, output_path.name, tmpdir)

        if tool_id == "pdf-to-powerpoint":
            output_path = tmpdir / f"{stem}.pptx"
            _pdf_to_powerpoint(input_path, output_path, tmpdir)
            return _as_download(output_path, output_path.name, tmpdir)

        if tool_id == "pdf-to-excel":
            output_path = tmpdir / f"{stem}.xlsx"
            _pdf_to_excel(input_path, output_path)
            return _as_download(output_path, output_path.name, tmpdir)

        if tool_id in {"word-to-pdf", "powerpoint-to-pdf", "excel-to-pdf"}:
            output_path = _libreoffice_convert(input_path, "pdf", tmpdir)
            return _as_download(output_path, output_path.name, tmpdir)

        if tool_id == "unlock-pdf":
            output_path = tmpdir / f"{stem}_unlocked.pdf"
            with pikepdf.open(str(input_path), password=password or None) as pdf:
                pdf.save(str(output_path))
            return _as_download(output_path, output_path.name, tmpdir)

        if tool_id == "protect-pdf":
            if not password:
                raise HTTPException(status_code=400, detail="Password is required to protect a PDF.")

            output_path = tmpdir / f"{stem}_protected.pdf"
            reader = PdfReader(str(input_path))
            writer = PdfWriter()

            for page in reader.pages:
                writer.add_page(page)

            writer.encrypt(
                user_password=password,
                owner_password=owner_password or password,
                use_128bit=True,
            )

            with output_path.open("wb") as f:
                writer.write(f)

            return _as_download(output_path, output_path.name, tmpdir)

        if tool_id == "pdf-to-pdfa":
            output_path = tmpdir / f"{stem}_pdfa.pdf"
            _convert_to_pdfa(input_path, output_path)
            return _as_download(output_path, output_path.name, tmpdir)

        if tool_id == "repair-pdf":
            output_path = tmpdir / f"{stem}_repaired.pdf"
            with pikepdf.open(str(input_path), attempt_recovery=True) as pdf:
                pdf.save(str(output_path))
            return _as_download(output_path, output_path.name, tmpdir)

        raise HTTPException(status_code=404, detail="Unknown tool ID.")

    except HTTPException:
        shutil.rmtree(tmpdir, ignore_errors=True)
        raise
    except Exception as exc:  # pragma: no cover
        logger.exception("Unexpected tool processing error for tool_id=%s", tool_id)
        shutil.rmtree(tmpdir, ignore_errors=True)
        raise HTTPException(status_code=500, detail="Failed to process file.") from exc


@app.exception_handler(HTTPException)
async def http_exc_handler(_, exc: HTTPException):
    return JSONResponse(status_code=exc.status_code, content={"error": exc.detail})
